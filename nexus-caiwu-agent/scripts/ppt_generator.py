#!/usr/bin/env python3
"""
本地PPT生成模块 - 增强版
使用 python-pptx 和 matplotlib 生成带图表的财报分析PPT
"""

from typing import Dict, Any, Optional, List
from pathlib import Path
import sys
import io

# 尝试导入 python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# 尝试导入 matplotlib
try:
    import matplotlib
    matplotlib.use('Agg')  # 使用非交互式后端
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    from matplotlib.patches import Circle, Rectangle, FancyBboxPatch
    import numpy as np

    # 配置中文字体 - 优先使用 Windows 系统字体
    plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei', 'KaiTi', 'FangSong', 'SimSun']
    plt.rcParams['axes.unicode_minus'] = False  # 解决负号显示问题
    plt.rcParams['font.size'] = 10

    HAS_MATPLOTLIB = True
except ImportError:
    HAS_MATPLOTLIB = False

# 图表配色方案 - 与HTML报告一致
COLORS = {
    'primary': '#4CAF50',      # 绿色
    'success': '#66BB6A',      # 浅绿
    'warning': '#FFA726',      # 橙色
    'danger': '#EF5350',       # 红色
    'info': '#42A5F5',         # 蓝色
    'purple': '#AB47BC',       # 紫色
    'cyan': '#26C6DA',         # 青色
    'background': '#F5F5F5',   # 浅灰背景
    'text': '#37474F',         # 深灰文字
    'grid': '#E0E0E0'          # 网格线
}


class LocalPPTGenerator:
    """本地PPT生成器"""

    def __init__(self):
        if not HAS_PPTX:
            raise ImportError("需要安装 python-pptx 库: pip install python-pptx")

    def generate_financial_report(
        self,
        data: Dict[str, Any],
        output_path: str,
        style: str = "business"
    ) -> str:
        """
        生成财务分析PPT - 增强版（带图表）

        Args:
            data: 财务分析数据
            output_path: 输出文件路径
            style: PPT风格

        Returns:
            生成的PPT文件路径
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # 获取基本信息
        stock_name = data.get("stock_name", "未知公司")
        stock_code = data.get("stock_code", "")
        health_score = data.get("health_score", "N/A")
        risk_level = data.get("risk_level", "N/A")
        key_metrics = data.get("key_metrics", {})
        health_details = data.get("health_details", {})

        # 获取行业对比数据
        industry_comparison = data.get("industry_comparison", {})

        # 第1页：封面
        self._add_cover_slide(prs, stock_name, stock_code)

        # 第2页：公司概况
        self._add_company_overview(prs, stock_name, stock_code, data)

        # 第3页：核心财务指标
        self._add_key_metrics(prs, stock_name, key_metrics)

        # 第4页：健康评分（带仪表盘和雷达图）
        self._add_health_score_with_charts(prs, health_score, risk_level, health_details)

        # 第5页：行业对比分析（带条形图）
        if industry_comparison:
            self._add_industry_comparison(prs, key_metrics, industry_comparison)

        # 第6页：盈利能力分析
        self._add_profitability_analysis(prs, key_metrics)

        # 第7页：投资建议
        self._add_recommendations(prs, data.get("analysis", {}).get("recommendations", []))

        # 第8页：风险提示
        self._add_disclaimer(prs)

        # 保存PPT
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        return str(output_path)

    def _add_cover_slide(self, prs: Presentation, stock_name: str, stock_code: str):
        """添加封面"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 使用标题布局的幻灯片
        # 标题 - 手动创建文本框
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(1.5)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = f"{stock_name}财务分析报告"
        text_frame.paragraphs[0].font.size = Pt(36)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 副标题
        top = Inches(3)
        height = Inches(1)

        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = subtitle_box.text_frame
        text_frame.text = f"股票代码：{stock_code}"
        text_frame.paragraphs[0].font.size = Pt(24)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 日期
        from datetime import datetime
        top = Inches(4.5)
        date_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = date_box.text_frame
        text_frame.text = f"报告日期：{datetime.now().strftime('%Y年%m月%d日')}"
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_company_overview(self, prs: Presentation, stock_name: str, stock_code: str, data: Dict):
        """添加公司概况"""
        self._add_title_slide(prs, "公司概况")

        # 添加内容
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # 基本信息
        p = text_frame.paragraphs[0]
        p.text = "基本信息"
        p.font.bold = True
        p.font.size = Pt(20)

        p = text_frame.add_paragraph()
        p.text = f"公司名称：{stock_name}"
        p.font.size = Pt(16)
        p.level = 0

        p = text_frame.add_paragraph()
        p.text = f"股票代码：{stock_code}"
        p.font.size = Pt(16)
        p.level = 0

        # 安全获取行业名称
        industry_data = data.get("industry", "N/A")
        if isinstance(industry_data, dict):
            industry = industry_data.get("name", "N/A")
        else:
            industry = str(industry_data) if industry_data else "N/A"

        p = text_frame.add_paragraph()
        p.text = f"所属行业：{industry}"
        p.font.size = Pt(16)
        p.level = 0

    def _add_key_metrics(self, prs: Presentation, stock_name: str, metrics: Dict):
        """添加核心财务指标"""
        self._add_title_slide(prs, "核心财务指标")

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(4)
        height = Inches(5)

        # 左侧：营收与利润
        textbox1 = slide.shapes.add_textbox(left, top, width, height)
        tf1 = textbox1.text_frame

        p = tf1.paragraphs[0]
        p.text = "营收与利润"
        p.font.bold = True
        p.font.size = Pt(18)

        revenue = metrics.get("revenue_billion", "N/A")
        p = tf1.add_paragraph()
        p.text = f"营业收入：{revenue}亿元"
        p.font.size = Pt(14)

        profit = metrics.get("net_profit_billion", "N/A")
        p = tf1.add_paragraph()
        p.text = f"净利润：{profit}亿元"
        p.font.size = Pt(14)

        margin = metrics.get("net_profit_margin", "N/A")
        p = tf1.add_paragraph()
        p.text = f"净利率：{margin}%"
        p.font.size = Pt(14)

        # 右侧：盈利能力
        left = Inches(5)
        textbox2 = slide.shapes.add_textbox(left, top, width, height)
        tf2 = textbox2.text_frame

        p = tf2.paragraphs[0]
        p.text = "盈利能力"
        p.font.bold = True
        p.font.size = Pt(18)

        gross_margin = metrics.get("gross_margin", "N/A")
        p = tf2.add_paragraph()
        p.text = f"毛利率：{gross_margin}%"
        p.font.size = Pt(14)

        roe = metrics.get("roe", "N/A")
        p = tf2.add_paragraph()
        p.text = f"ROE：{roe}%"
        p.font.size = Pt(14)

        roa = metrics.get("roa", "N/A")
        p = tf2.add_paragraph()
        p.text = f"ROA：{roa}%"
        p.font.size = Pt(14)

    def _add_health_score(self, prs: Presentation, score: Any, risk_level: str, details: Dict):
        """添加健康评分"""
        self._add_title_slide(prs, "健康评分")

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 总体评分
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(1)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = f"健康评分：{score}/100"
        p.font.bold = True
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

        # 风险等级
        p = tf.add_paragraph()
        p.text = f"风险等级：{risk_level}"
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER

        # 各维度评分
        top = Inches(4)
        textbox2 = slide.shapes.add_textbox(left, top, width, height)
        tf2 = textbox2.text_frame

        p = tf2.paragraphs[0]
        p.text = "分维度评分："
        p.font.bold = True
        p.font.size = Pt(16)

        dimension_names = {
            "profitability": "盈利能力",
            "solvency": "偿债能力",
            "efficiency": "运营效率",
            "growth": "成长能力",
            "cashflow": "现金流质量"
        }

        for dim, detail in details.items():
            if isinstance(detail, dict) and "score" in detail:
                name = dimension_names.get(dim, dim)
                score_val = detail.get("score", "N/A")
                max_val = detail.get("max", 25)
                p = tf2.add_paragraph()
                p.text = f"{name}：{score_val}/{max_val}"
                p.font.size = Pt(14)
                p.level = 1

    def _add_health_score_with_charts(self, prs: Presentation, score: Any, risk_level: str, details: Dict):
        """添加健康评分（带图表）"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 标题：健康评分
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.6)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"健康评分：{score}/100  |  风险等级：{risk_level}"
        p.font.bold = True
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER

        # 左侧：仪表盘图
        if HAS_MATPLOTLIB and isinstance(score, (int, float)):
            gauge_chart = self._create_score_gauge_chart(float(score))
            if gauge_chart:
                self._add_picture_to_slide(slide, gauge_chart, 0.5, 1.2, 4, 4)

        # 右侧：雷达图
        if HAS_MATPLOTLIB and details:
            radar_chart = self._create_radar_chart(details)
            if radar_chart:
                self._add_picture_to_slide(slide, radar_chart, 5, 1.2, 4.5, 4)

    def _add_industry_comparison(self, prs: Presentation, metrics: Dict, industry_comparison: Dict):
        """添加行业对比分析（带条形图）"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 标题
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.6)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "行业对比分析"
        p.font.bold = True
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

        # 添加条形图
        if HAS_MATPLOTLIB and industry_comparison:
            bar_chart = self._create_bar_chart(metrics, industry_comparison)
            if bar_chart:
                self._add_picture_to_slide(slide, bar_chart, 1, 1.5, 8, 4.5)

        # 添加说明文字
        left = Inches(1)
        top = Inches(6.2)
        width = Inches(8)
        height = Inches(0.6)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "注：绿色柱代表本公司数据，蓝色柱代表行业平均水平"
        p.font.size = Pt(11)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    def _add_profitability_analysis(self, prs: Presentation, metrics: Dict):
        """添加盈利能力分析"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 添加标题
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.6)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "盈利能力分析"
        p.font.bold = True
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

        # 内容区域
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        # 净利率分析
        net_margin = metrics.get("net_profit_margin", 0)
        p = tf.paragraphs[0]
        p.text = f"净利率：{net_margin}%"
        p.font.size = Pt(18)

        if net_margin >= 50:
            assessment = "净利率超过50%，盈利能力极强，具有强大的定价能力和成本控制能力。"
        elif net_margin >= 20:
            assessment = "净利率处于20%-50%区间，盈利能力良好。"
        else:
            assessment = "净利率低于20%，建议关注成本控制和产品定价能力。"

        p = tf.add_paragraph()
        p.text = assessment
        p.font.size = Pt(14)

        # ROE分析
        p = tf.add_paragraph()
        p.text = ""
        p = tf.add_paragraph()
        roe = metrics.get("roe", 0)
        p.text = f"ROE：{roe}%"
        p.font.size = Pt(18)

        if roe >= 20:
            assessment = "ROE超过20%，股东回报率极高，资本利用效率优秀。"
        elif roe >= 10:
            assessment = "ROE处于10%-20%区间，股东回报率良好。"
        else:
            assessment = "ROE低于10%，建议提升资产使用效率。"

        p = tf.add_paragraph()
        p.text = assessment
        p.font.size = Pt(14)

    def _add_recommendations(self, prs: Presentation, recommendations: list):
        """添加投资建议"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 添加标题
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.6)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "投资建议"
        p.font.bold = True
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

        # 内容区域
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(5.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.paragraphs[0]
        p.text = "综合评估："
        p.font.bold = True
        p.font.size = Pt(18)

        for i, rec in enumerate(recommendations, 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {rec}"
            p.font.size = Pt(14)
            p.level = 1

    def _add_disclaimer(self, prs: Presentation):
        """添加风险提示"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 添加标题
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.6)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "风险提示"
        p.font.bold = True
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

        # 内容区域
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(5.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        disclaimer_text = """本报告基于公开财务数据进行分析，仅供参考。

投资风险提示：
• 股票投资存在市场风险
• 过往业绩不代表未来表现
• 请结合个人风险承受能力投资

报告生成：nexus-caiwu-agent
数据来源：公开财务数据
免责声明：本报告不构成投资建议"""

        tf.text = disclaimer_text
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(12)

    # ===== 图表生成方法 =====

    def _create_score_gauge_chart(self, score: float, max_score: float = 100) -> Optional[bytes]:
        """创建健康评分仪表盘图表"""
        if not HAS_MATPLOTLIB:
            return None

        try:
            fig, ax = plt.subplots(figsize=(4, 4))
            fig.patch.set_facecolor('white')
            ax.set_facecolor('white')

            # 计算角度
            percentage = min(score / max_score, 1)
            angles = np.linspace(0, np.pi, 100)

            # 背景弧
            ax.plot(np.cos(angles), np.sin(angles), color='#E0E0E0', linewidth=20)

            # 进度弧 - 使用渐变色效果
            if score >= 60:
                color = COLORS['primary']  # 绿色
            elif score >= 40:
                color = COLORS['warning']  # 橙色
            else:
                color = COLORS['danger']   # 红色

            angles_filled = np.linspace(0, np.pi * percentage, int(100 * percentage))
            ax.plot(np.cos(angles_filled), np.sin(angles_filled), color=color, linewidth=20)

            # 中心分数
            ax.text(0, -0.15, f'{score:.0f}', fontsize=48, fontweight='bold',
                   ha='center', va='center', color=COLORS['text'])
            ax.text(0, -0.35, f'/{max_score:.0f}', fontsize=24, ha='center', va='center',
                   color=COLORS['text'], alpha=0.6)

            ax.set_xlim(-1.3, 1.3)
            ax.set_ylim(-0.5, 1.3)
            ax.axis('off')

            # 保存到内存
            buf = io.BytesIO()
            plt.tight_layout(pad=0)
            plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            buf.seek(0)
            plt.close(fig)

            return buf.read()
        except Exception as e:
            print(f"[警告] 创建仪表盘图表失败: {e}")
            plt.close('all')
            return None

    def _create_radar_chart(self, health_details: Dict[str, Any]) -> Optional[bytes]:
        """创建雷达图"""
        if not HAS_MATPLOTLIB:
            return None

        try:
            # 提取维度数据
            dimensions = []
            scores = []
            max_scores = []

            dimension_names = {
                "profitability": "盈利能力",
                "solvency": "偿债能力",
                "efficiency": "运营效率",
                "growth": "成长能力",
                "cashflow": "现金流质量"
            }

            for dim, detail in health_details.items():
                if isinstance(detail, dict) and "score" in detail:
                    name = dimension_names.get(dim, dim)
                    score = detail.get("score", 0)
                    max_val = detail.get("max", 25)
                    dimensions.append(name)
                    scores.append(score)
                    max_scores.append(max_val)

            if len(dimensions) < 3:
                return None

            # 计算归一化分数 (0-100)
            normalized_scores = [(s / m * 100) if m > 0 else 0 for s, m in zip(scores, max_scores)]

            # 创建雷达图
            categories = dimensions
            N = len(categories)

            # 计算角度
            angles = [n / float(N) * 2 * np.pi for n in range(N)]
            angles += angles[:1]
            normalized_scores += normalized_scores[:1]

            fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(projection='polar'))
            fig.patch.set_facecolor('white')

            # 绘制雷达图
            ax.plot(angles, normalized_scores, 'o-', linewidth=2, color=COLORS['primary'], label='实际得分')
            ax.fill(angles, normalized_scores, alpha=0.25, color=COLORS['primary'])

            # 添加参考线 (满分)
            max_values = [100] * N + [100]
            ax.plot(angles, max_values, '--', linewidth=1, color=COLORS['grid'], alpha=0.5, label='满分')

            # 设置刻度
            ax.set_xticks(angles[:-1])
            ax.set_xticklabels(categories, fontsize=11, color=COLORS['text'])
            ax.set_ylim(0, 100)
            ax.set_yticks([20, 40, 60, 80, 100])
            ax.set_yticklabels(['20', '40', '60', '80', '100'], fontsize=9, color=COLORS['text'], alpha=0.6)
            ax.grid(True, color=COLORS['grid'], alpha=0.3)

            # 添加背景色
            ax.set_facecolor('#FAFAFA')

            # 保存到内存
            buf = io.BytesIO()
            plt.tight_layout(pad=0.5)
            plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            buf.seek(0)
            plt.close(fig)

            return buf.read()
        except Exception as e:
            print(f"[警告] 创建雷达图失败: {e}")
            plt.close('all')
            return None

    def _create_bar_chart(self, metrics: Dict[str, Any], industry_comparison: Dict[str, Any]) -> Optional[bytes]:
        """创建行业对比条形图"""
        if not HAS_MATPLOTLIB:
            return None

        try:
            # 提取对比数据
            metrics_to_show = []
            company_values = []
            industry_values = []

            # 选择几个关键指标
            key_metrics = [
                ("gross_margin", "毛利率"),
                ("net_margin", "净利率"),
                ("roe", "ROE"),
                ("debt_ratio", "资产负债率")
            ]

            for metric_key, metric_name in key_metrics:
                if metric_key in industry_comparison:
                    comp = industry_comparison[metric_key]
                    if isinstance(comp, dict):
                        company_val = comp.get("company_value", 0)
                        industry_val = comp.get("industry_ideal", 0)
                        metrics_to_show.append(metric_name)
                        company_values.append(company_val)
                        industry_values.append(industry_val)

            if not metrics_to_show:
                return None

            # 创建条形图
            fig, ax = plt.subplots(figsize=(8, 5))
            fig.patch.set_facecolor('white')

            x = np.arange(len(metrics_to_show))
            width = 0.35

            bars1 = ax.bar(x - width/2, company_values, width, label='本公司', color=COLORS['primary'], alpha=0.8)
            bars2 = ax.bar(x + width/2, industry_values, width, label='行业平均', color=COLORS['info'], alpha=0.8)

            ax.set_xlabel('财务指标', fontsize=12, color=COLORS['text'])
            ax.set_ylabel('数值 (%)', fontsize=12, color=COLORS['text'])
            ax.set_xticks(x)
            ax.set_xticklabels(metrics_to_show, fontsize=11, color=COLORS['text'])
            ax.legend(fontsize=11, loc='upper right')
            ax.set_facecolor('#FAFAFA')
            ax.grid(True, axis='y', color=COLORS['grid'], alpha=0.3)

            # 添加数值标签
            for bars in [bars1, bars2]:
                for bar in bars:
                    height = bar.get_height()
                    ax.annotate(f'{height:.1f}%',
                              xy=(bar.get_x() + bar.get_width() / 2, height),
                              xytext=(0, 3),
                              textcoords="offset points",
                              ha='center', va='bottom',
                              fontsize=9, color=COLORS['text'])

            plt.tight_layout()

            # 保存到内存
            buf = io.BytesIO()
            plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            buf.seek(0)
            plt.close(fig)

            return buf.read()
        except Exception as e:
            print(f"[警告] 创建条形图失败: {e}")
            plt.close('all')
            return None

    def _add_picture_to_slide(self, slide, picture_bytes: bytes, left: float, top: float,
                              width: float, height: float):
        """将图片添加到幻灯片"""
        try:
            # 将字节数据保存到临时文件
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                tmp.write(picture_bytes)
                tmp_path = tmp.name

            # 添加图片
            slide.shapes.add_picture(tmp_path, Inches(left), Inches(top),
                                     width=Inches(width), height=Inches(height))

            # 删除临时文件
            import os
            try:
                os.unlink(tmp_path)
            except:
                pass
        except Exception as e:
            print(f"[警告] 添加图片到幻灯片失败: {e}")

    def _add_title_slide(self, prs: Presentation, title: str):
        """添加标题页"""
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        # 尝试使用标题形状，如果不存在则创建文本框
        try:
            title_shape = slide.shapes.title
            if title_shape is None:
                raise AttributeError("No title shape in layout")
            title_shape.text = title
        except (AttributeError, KeyError):
            # 手动创建标题文本框
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)

            title_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = title_box.text_frame
            text_frame.text = title
            text_frame.paragraphs[0].font.size = Pt(32)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def create_ppt_generator() -> LocalPPTGenerator:
    """创建PPT生成器"""
    return LocalPPTGenerator()


# 测试代码
if __name__ == "__main__":
    if not HAS_PPTX:
        print("需要安装 python-pptx: pip install python-pptx")
        sys.exit(1)

    # 测试数据
    test_data = {
        "stock_name": "贵州茅台",
        "stock_code": "600519",
        "health_score": 82,
        "risk_level": "低风险",
        "key_metrics": {
            "revenue_billion": 1309.04,
            "net_profit_billion": 668.99,
            "net_profit_margin": 51.11,
            "gross_margin": 68.35,
            "roe": 25.18,
            "roa": 21.95
        },
        "health_details": {
            "profitability": {"score": 25, "max": 25},
            "solvency": {"score": 25, "max": 25},
            "efficiency": {"score": 5, "max": 20},
            "growth": {"score": 12, "max": 15},
            "cashflow": {"score": 15, "max": 15}
        },
        "analysis": {
            "recommendations": [
                "财务状况优异",
                "盈利能力突出",
                "建议关注"
            ]
        }
    }

    # 生成PPT
    generator = create_ppt_generator()
    output_path = generator.generate_financial_report(
        test_data,
        "test_report.pptx"
    )

    print(f"PPT已生成: {output_path}")
